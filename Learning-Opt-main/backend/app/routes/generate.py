from flask import Blueprint, request, jsonify, send_file
from flask_cors import cross_origin
from pptx import Presentation
from openpyxl import load_workbook
import os
import tempfile
import json
import io
import re
from copy import deepcopy
from datetime import datetime
import logging
from openpyxl.utils import get_column_letter

bp = Blueprint('generate', __name__, url_prefix='/generate')

# Base paths
BASE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
TEMPLATE_DIR = os.path.join(BASE_DIR, 'uploads', 'templates')
OUTPUT_DIR = os.path.join(BASE_DIR, 'static', 'generated')
os.makedirs(OUTPUT_DIR, exist_ok=True)

def to_number(val):
    """Convert to int/float if numeric, else return original or None."""
    try:
        if val is None or str(val).strip() == "":
            return None
        num = float(val)
        return int(num) if num.is_integer() else num
    except (ValueError, TypeError):
        return val  # Keep text as-is

# --------- Excel Generation route (/generate/excel) ---------

# Unmerge ALL merged cells before writing

# Now write without extra merged-cell checks:
def safe_write(ws, row, col_idx, value):
    ws.cell(row=row, column=col_idx).value = value


@bp.route('/excel', methods=['POST'])
def generate_excel():
    try:
        students = request.json.get("students", [])
        if not students:
            return jsonify({"error": "No student data received"}), 400

        template_path = os.path.join(TEMPLATE_DIR, 'grades2.xlsx')
        if not os.path.exists(template_path):
            return jsonify({"error": "Grades.xlsx template not found"}), 500

        wb = load_workbook(template_path)
        sheet_map = {}
        for dept in ["PRODUCTION", "SUPPORT", "TECHNICAL"]:
            if dept in wb.sheetnames:
                sheet_map[dept] = wb[dept]

        # --- Unmerge all merged cells from row 10 onwards to avoid merged cell write errors ---
        for ws in sheet_map.values():
            merged_ranges = list(ws.merged_cells.ranges)
            for merged_range in merged_ranges:
                ws.unmerge_cells(str(merged_range))

        first_student = students[0]
        immersion_date = first_student.get("date_of_immersion", "")
        batch = first_student.get("batch", "")
        school = first_student.get("school", "")
    
        for dept, ws in sheet_map.items():
            safe_write(ws, 8, 8, f"{batch} - {school}")  # H8
            safe_write(ws, 9, 8, immersion_date)          # H9

        row_counter = {dept: 10 for dept in sheet_map.keys()}

        for s in students:
            dept_raw = (s.get("department") or "").strip().upper()
            if dept_raw in ["TECHNICAL", "IT"]:
                dept = "TECHNICAL"
            elif dept_raw == "PROD":
                dept = "PRODUCTION"
            else:
                dept = "SUPPORT"

            if dept not in sheet_map:
                logging.warning(f"Department {dept} not in sheet_map, skipping student: {s}")
                continue

            ws = sheet_map[dept]
            row = row_counter[dept]

            # Write student info columns B-F (2-6)
            safe_write(ws, row, 2, s.get("last_name", ""))
            safe_write(ws, row, 3, s.get("first_name", ""))
            safe_write(ws, row, 4, s.get("middle_name", ""))
            safe_write(ws, row, 5, s.get("strand", ""))
            safe_write(ws, row, 6, s.get("department", ""))

            # Grades columns G-R (7-18) = 1G to 12G
            for i, col_idx in enumerate(range(7, 19), start=1):
                val = to_number(s.get(f"{i}G", ""))
                safe_write(ws, row, col_idx, val)
                logging.debug(f"Wrote {val} to {get_column_letter(col_idx)}{row}")

            # Extra columns per department
            if dept == "PRODUCTION":
                extras = {
                    22: "13G",  # V
                    23: "14G",  # W
                    24: "15G",  # X
                    25: "16G",  # Y
                    28: "17G",  # AB
                    29: "18G",  # AC
                }
            elif dept == "SUPPORT":
                extras = {
                    21: "13G",  # U
                    26: "14G",  # Z
                    29: "15G",  # AC
                }
            elif dept == "TECHNICAL":
                extras = {
                    20: "13G",  # T
                    27: "14G",  # AA
                    29: "15G",  # AC
                }
            else:
                extras = {}

            for col_idx, key in extras.items():
                val = to_number(s.get(key, ""))
                safe_write(ws, row, col_idx, val)
                logging.debug(f"Wrote {val} to {get_column_letter(col_idx)}{row}")

            row_counter[dept] += 1

        # Save locally for manual checking
        debug_path = os.path.join(OUTPUT_DIR, "debug_generated.xlsx")
        wb.save(debug_path)
        logging.info(f"Saved debug Excel file to: {debug_path}")

        # Save temp file for sending
        temp_dir = tempfile.mkdtemp()
        output_filename = "generated_immersion_report.xlsx"
        output_path = os.path.join(temp_dir, output_filename)
        wb.save(output_path)

        return send_file(
            output_path,
            as_attachment=True,
            download_name=output_filename
        )

    except Exception as e:
        logging.error(f"Error generating Excel: {e}")
        return jsonify({"error": str(e)}), 500

# --------- Certificate generation and preview routes ---------

def fill_slide(slide, data_row):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            full_text = ''.join(run.text for run in paragraph.runs)
            replaced = full_text
            for key, value in data_row.items():
                ph = f"{{{key}}}"
                replaced = replaced.replace(ph, str(value))
            if replaced != full_text:
                for run in paragraph.runs:
                    run.text = ""
                paragraph.runs[0].text = replaced

@bp.route('/delete_certificate', methods=['DELETE'])
@cross_origin()
def delete_certificate():
    filename = request.args.get("filename")
    if not filename:
        return jsonify({"error": "Filename is required"}), 400

    file_path = os.path.join(OUTPUT_DIR, filename)
    if not os.path.exists(file_path):
        return jsonify({"error": "File not found"}), 404

    try:
        os.remove(file_path)
        return jsonify({"message": "File deleted successfully"}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@bp.route('/certificates', methods=['POST', 'OPTIONS'])
@cross_origin()
def generate_certificates():
    data = request.get_json()
    template_type = data.get('template', 'ojt')
    rows = data.get('rows', [])

    if not rows:
        return jsonify({"error": "No data provided"}), 400

    tpl_filename = f"{template_type}.pptx"
    template_path = os.path.join(TEMPLATE_DIR, tpl_filename)
    if not os.path.exists(template_path):
        return jsonify({"error": f"Template '{tpl_filename}' not found"}), 404

    prs = Presentation(template_path)
    source_slide = prs.slides[0]
    original_elements = [deepcopy(shape.element) for shape in source_slide.shapes]
    fill_slide(source_slide, rows[0])

    for row in rows[1:]:
        new_slide = prs.slides.add_slide(source_slide.slide_layout)
        for shp in list(new_slide.shapes):
            new_slide.shapes._spTree.remove(shp.element)
        for el in original_elements:
            new_slide.shapes._spTree.append(deepcopy(el))
        fill_slide(new_slide, row)

    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M")
    output_name = f"certificate_{template_type} ({timestamp}).pptx"
    output_path = os.path.join(OUTPUT_DIR, output_name)
    prs.save(output_path)

    return jsonify({"message": "Certificates generated", "files": [output_name]})

@bp.route('/files/<filename>', methods=['GET'])
@cross_origin()
def get_generated_file(filename):
    file_path = os.path.join(OUTPUT_DIR, filename)
    if not os.path.exists(file_path):
        return jsonify({"error": "File not found"}), 404
    return send_file(file_path, as_attachment=True, download_name=filename)

@bp.route('/preview', methods=['POST', 'OPTIONS'])
@cross_origin()
def preview_certificate():
    data = request.get_json()
    template_type = data.get('template', 'ojt')
    rows = data.get('rows', [])

    if not rows:
        return jsonify({"error": "No data to preview"}), 400

    tpl_filename = f"{template_type}.pptx"
    template_path = os.path.join(TEMPLATE_DIR, tpl_filename)

    if not os.path.exists(template_path):
        return jsonify({"error": f"Template '{tpl_filename}' not found"}), 404

    if template_type == 'tesda':
        # TESDA preview HTML generation logic here (same as before)...
        # (I can help add this if you want, but omitted here for brevity)
        pass

    # Default preview (non-TESDA)
    html_parts = [
        "<!DOCTYPE html><html><head><meta charset='utf-8'>",
        "<meta name='viewport' content='width=device-width, initial-scale=1.0'>",
        "<title>Certificate Preview</title><style>",
        "body { font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; background: #2d3748; margin: 0; padding: 20px; min-height: 100vh; }",
        ".container { max-width: 900px; margin: 0 auto; }",
        "h2 { text-align: center; color: white; font-size: 2.5rem; margin-bottom: 30px; }",
        ".slide-preview { background: #4a5568; border-radius: 12px; padding: 25px; margin-bottom: 25px; box-shadow: 0 8px 32px rgba(0,0,0,0.2); border: 1px solid #718096; transition: transform 0.2s ease; }",
        ".slide-preview:hover { transform: translateY(-2px); box-shadow: 0 12px 40px rgba(0,0,0,0.3); }",
        ".slide-preview h4 { color: white; font-size: 1.3rem; margin: 0 0 15px 0; padding-bottom: 8px; border-bottom: 2px solid #718096; }",
        ".slide-preview p { margin: 10px 0; font-size: 1.1rem; line-height: 1.6; color: #e2e8f0; padding: 8px 12px; background: #2d3748; border-radius: 6px; border-left: 4px solid #a361ef; }",
        ".certificate-text { font-weight: 500; }",
        "@media (max-width: 768px) { .container { padding: 10px; } .slide-preview { padding: 20px; } h2 { font-size: 2rem; } }",
        "</style></head><body><div class='container'><h2>Certificate Preview</h2>"
    ]

    for idx, row in enumerate(rows):
        prs_row = Presentation(template_path)
        slide = prs_row.slides[0]
        fill_slide(slide, row)
        html_parts.append("<div class='slide-preview certificate-text'>")
        html_parts.append(f"<h4>Certificate {idx+1}</h4>")

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = ''.join(run.text for run in paragraph.runs)
                if text.strip():
                    html_parts.append(f"<p>{text}</p>")
        html_parts.append("</div>")
    html_parts.append("</div></body></html>")
    return "\n".join(html_parts), 200, {"Content-Type": "text/html"}
