from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from openpyxl import load_workbook
from datetime import datetime
import os
import uuid
import io
import json
import re
import traceback
import requests  # <-- needed for internal HTTP calls

from app.routes.auth import auth_bp
from app.routes.generate import bp as generate_bp
from app.routes.upload import upload_bp          # <-- fixed
from app.routes.excel_generate import excel_bp
from app.routes.immersion import immersion_bp

from app import config

app = Flask(__name__)
CORS(app)

CORS(app, resources={r"/api/*": {"origins": "*"}})

@app.after_request
def expose_headers(resp):
    resp.headers["Access-Control-Expose-Headers"] = "Content-Disposition"
    return resp

BASE_DIR = os.path.dirname(__file__)
UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads", "templates")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

GENERATED_FOLDER = os.path.join("static", "generated")
os.makedirs(GENERATED_FOLDER, exist_ok=True)

# Blueprints
app.register_blueprint(auth_bp)
app.register_blueprint(generate_bp)
app.register_blueprint(upload_bp)
app.register_blueprint(excel_bp)
app.register_blueprint(immersion_bp)

PLACEHOLDER_RE = re.compile(r"\{([^}]+)\}")
recent_downloads = []

# ---------- Utilities ----------

def format_value(val, fmt=None):
    return "" if val is None else str(val)

def replace_placeholders_in_cell(text, mapping, rowdict):
    if "YEAR LAST ATTENDED" in text.upper():
        context = None
        up = text.upper()
        if "ELEMENTARY" in up:
            context = "ELEMENTARY"
        elif "SECONDARY" in up:
            context = "SECONDARY"
        elif "TERTIARY" in up:
            context = "TERTIARY"
    else:
        context = None

    def repl(m):
        key = m.group(1)
        mp = mapping.get(key, key)
        if isinstance(mp, dict):
            col = mp.get(context) or mp.get("DEFAULT")
        else:
            col = mp
        val = rowdict.get(col, "")
        return format_value(val)

    return PLACEHOLDER_RE.sub(repl, text)

def replace_placeholders_in_worksheet(ws, mapping, rowdict):
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
        for cell in row:
            if isinstance(cell.value, str) and "{" in cell.value and "}" in cell.value:
                cell.value = replace_placeholders_in_cell(cell.value, mapping, rowdict)

def _safe_sheet_title(s: str, used: set) -> str:
    title = (s or "").strip() or "Row"
    for ch in '[]:*?/\\':
        title = title.replace(ch, "-")
    title = title[:31] or "Row"
    orig = title
    i = 2
    while title in used:
        suffix = f" ({i})"
        title = (orig[: 31 - len(suffix)] + suffix) if len(orig) + len(suffix) > 31 else orig + suffix
        i += 1
    used.add(title)
    return title

def _copy_template_sheet_with_fallback(wb, template_ws, new_title):
    try:
        ws_copy = wb.copy_worksheet(template_ws)
        ws_copy.title = new_title
        return ws_copy
    except Exception as e:
        print("[WARN] copy_worksheet failed; falling back:", repr(e))
        ws = wb.create_sheet(title=new_title)
        for rng in template_ws.merged_cells.ranges:
            ws.merge_cells(str(rng))
        for r in range(1, template_ws.max_row + 1):
            for c in range(1, template_ws.max_column + 1):
                v = template_ws.cell(row=r, column=c).value
                if v is not None:
                    ws.cell(row=r, column=c, value=v)
        return ws

def is_top_left_merged_cell(ws, row, col):
    cell = ws.cell(row=row, column=col)
    for merged_range in ws.merged_cells.ranges:
        if cell.coordinate in merged_range:
            return cell.coordinate == merged_range.start_cell.coordinate
    return True

def to_number(val):
    try:
        return int(val)
    except (ValueError, TypeError):
        try:
            return float(val)
        except (ValueError, TypeError):
            return val

def get_student_value(student, key):
    if not isinstance(student, dict):
        return None
    nk = key.strip().lower()

    for k, v in student.items():
        if isinstance(k, str) and k.strip().lower() == nk:
            return v

    for parent in ("scores", "grades", "appraisal", "performance"):
        pv = student.get(parent)
        if isinstance(pv, dict):
            val = get_student_value(pv, key)
            if val is not None:
                return val

    for k, v in student.items():
        if isinstance(k, str) and nk in k.strip().lower():
            return v
    return None

def force_full_calc_on_load(wb):
    try:
        wb.properties.calcPr.calcMode = "auto"
        wb.properties.calcPr.fullCalcOnLoad = True
        wb.calcPr.fullCalcOnLoad = True
    except Exception:
        pass

# ---------- Simple endpoints ----------

@app.route("/api/ping")
def ping():
    return jsonify(ok=True)

@app.route("/")
def home():
    return "Hello, Creo Certificate Backend!"

# ---------- Certificate generator ----------

@app.route('/generate/certificates', methods=['POST'])
def generate_certificates():
    data = request.json or {}
    template_path = data.get("templatePath")
    if not template_path or not os.path.exists(template_path):
        return jsonify({"error": "Invalid templatePath"}), 400

    output_folder = "static/generated"
    os.makedirs(output_folder, exist_ok=True)

    custom_filename = data.get("filename")
    if custom_filename:
        filename = f"{custom_filename}.pptx"
    else:
        name = data.get("name", "Certificate")
        filename = f"{name.replace(' ', '_')}_Certificate.pptx"

    output_path = os.path.join(output_folder, filename)

    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{{" in run.text and "}}" in run.text:
                            key = run.text.replace("{{", "").replace("}}", "").strip()
                            run.text = data.get(key, "")

    prs.save(output_path)
    return jsonify({"files": [filename]})

# ---------- TESDA upload → save copy ----------

@app.route('/api/generate', methods=['POST'])
def generate_tesda_excel():
    uploaded_file = request.files.get("file")
    if not uploaded_file:
        return jsonify({"error": "No file uploaded"}), 400

    temp_path = os.path.join(UPLOAD_FOLDER, f"temp_{uuid.uuid4().hex}.xlsx")
    uploaded_file.save(temp_path)

    try:
        wb = load_workbook(temp_path)
        now = datetime.now().strftime("%Y%m%d-%H%M%S")
        output_filename = f"tesda_record_{now}.xlsx"
        output_path = os.path.join(GENERATED_FOLDER, output_filename)
        wb.save(output_path)

        recent_downloads.insert(0, {
            "type": "tesda",
            "filename": output_filename,
            "timestamp": datetime.fromtimestamp(os.path.getmtime(output_path)).strftime("%Y-%m-%d %H:%M:%S"),
            "url": f"/static/generated/{output_filename}"
        })

        return send_file(output_path, as_attachment=True, download_name=output_filename)
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

# ---------- File listings (no duplicates) ----------

@app.route('/api/certificates', methods=['GET'])
def list_certificates():
    try:
        files = [f for f in os.listdir(GENERATED_FOLDER) if f.endswith(".pptx")]
        files.sort(key=lambda x: os.path.getmtime(os.path.join(GENERATED_FOLDER, x)), reverse=True)
        return jsonify(files)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/tesda', methods=['GET'])
def list_tesda():
    try:
        files = [f for f in os.listdir(GENERATED_FOLDER) if f.endswith(".xlsx")]
        files.sort(key=lambda x: os.path.getmtime(os.path.join(GENERATED_FOLDER, x)), reverse=True)
        return jsonify(files)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/download-history", methods=["GET"])
def get_download_history():
    folder = GENERATED_FOLDER
    files = [
        f for f in os.listdir(folder)
        if f.endswith(".pptx") or (f.endswith(".xlsx") and "tesda" in f.lower())
    ]
    files.sort(key=lambda f: os.path.getmtime(os.path.join(folder, f)), reverse=True)

    history = []
    for f in files:
        file_type = "certificate" if f.endswith(".pptx") else "tesda"
        history.append({
            "type": file_type,
            "filename": f,
            "timestamp": datetime.fromtimestamp(os.path.getmtime(os.path.join(folder, f))).strftime("%Y-%m-%d %H:%M"),
            "url": f"/static/generated/{f}"
        })
    return jsonify(history)

@app.route("/api/download-history", methods=["POST"])
def update_download_history():
    data = request.get_json() or {}
    filename = data.get("filename")
    if not filename:
        return jsonify({"error": "Missing filename"}), 400

    file_path = os.path.join(GENERATED_FOLDER, filename)
    if not os.path.exists(file_path):
        return jsonify({"error": "File does not exist"}), 404

    if not any(d.get("filename") == filename for d in recent_downloads):
        file_type = "tesda" if filename.lower().endswith(".xlsx") else "certificate"
        recent_downloads.insert(0, {
            "type": file_type,
            "filename": filename,
            "timestamp": datetime.fromtimestamp(os.path.getmtime(file_path)).strftime("%Y-%m-%d %H:%M"),
            "url": f"/static/generated/{filename}"
        })

    return jsonify({"success": True})

# ---------- Internal TESDA generator + proxy ----------

@app.route('/api/generate/excel', methods=['POST'])
def generate_excel_from_json():
    try:
        payload = request.get_json() or {}
        students = payload.get("students")
        if not students or not isinstance(students, list):
            return jsonify({"error": "Missing or invalid students data"}), 400

        template_path = os.path.join(app.config['UPLOAD_FOLDER'], "grades2.xlsx")
        if not os.path.exists(template_path):
            return jsonify({"error": f"Template not found: {template_path}"}), 404

        wb = load_workbook(template_path)

       

        sheet_mapping = {
            "PROD": "PRODUCTION",
            "IT": "TECHNICAL",
            "ACCTG": "SUPPORT",
            "ERT": "SUPPORT",
            "HS": "SUPPORT",
            "HSN": "SUPPORT",
            "ER": "SUPPORT"
        }
        start_rows = {sheet: 10 for sheet in ["PRODUCTION", "TECHNICAL", "SUPPORT"]}

        basic_mapping = {
            "last_name": 2,     # B
            "first_name": 3,    # C
            "middle_name": 4,   # D
            "strand": 5,        # E
            "department": 6,    # F
            "over_all": 7,      # G
            "total_score": 30   # AD (do not write directly)
        }

        score_mapping = {
            "wi": 8, "co": 9, "5s": 10, "bo": 11, "cbo": 12, "sdg": 13,
            "ohsa": 14, "we": 15, "ujc": 16, "iso": 17, "po": 18, "hr": 19,
            "perdev": 21, "supp": 26, "ds": 29
        }

        def has_name(stu):
            for key in ("last_name", "first_name", "name", "Name"):
                v = get_student_value(stu, key)
                if v and str(v).strip():
                    return True
            return False

        # Fill Excel with student data
        missing = []
        dept_students = {}
        for stu in students:
            if not has_name(stu):
                continue
            dept = (stu.get("department") or "").strip().upper()
            sheet_name = sheet_mapping.get(dept)
            if not sheet_name:
                continue
            dept_students.setdefault(sheet_name, []).append(stu)

        for sheet_name, stu_list in dept_students.items():
            ws = wb[sheet_name]
            row_num = start_rows[sheet_name]
            for stu in [s for s in stu_list if has_name(s)]:
                for key, col in basic_mapping.items():
                    if key == "total_score":
                        continue
                    if is_top_left_merged_cell(ws, row_num, col):
                        val = get_student_value(stu, key)
                        if key == "over_all":
                            val = to_number(val)
                        ws.cell(row=row_num, column=col, value=val or "")
                        if key == "over_all" and isinstance(val, (int, float)):
                            ws.cell(row=row_num, column=col).number_format = '0.0'

                for skey, col in score_mapping.items():
                    raw_val = get_student_value(stu, skey)
                    val = "" if raw_val is None else to_number(raw_val)
                    if raw_val is None:
                        missing.append({
                            "row_index": row_num,
                            "student": get_student_value(stu, "last_name") or get_student_value(stu, "first_name"),
                            "key": skey
                        })
                    if is_top_left_merged_cell(ws, row_num, col):
                        cell = ws.cell(row=row_num, column=col, value=val)
                        if isinstance(val, (int, float)):
                            cell.number_format = '0'
                row_num += 1
            start_rows[sheet_name] = row_num

        # --- Compute totals & grades ---
        written_fields = ["wi", "co", "5s", "bo", "cbo", "sdg"]
        performance_fields = ["ohsa", "we", "ujc", "iso", "po", "hr", "perdev", "supp", "ds"]

        for stu in students:
            for key in written_fields + performance_fields:
                val = get_student_value(stu, key) or 0
                try:
                    stu[key] = float(val)
                except ValueError:
                    stu[key] = 0.0

            total_score = sum(stu[k] for k in written_fields + performance_fields)
            stu["total_score"] = total_score

            stu["written_rating"] = round(sum(stu[k] for k in written_fields) / len(written_fields), 2)
            stu["performance_rating"] = round(sum(stu[k] for k in performance_fields) / len(performance_fields), 2)

            if total_score >= 90:
                stu["final_grade"] = "A"
            elif total_score >= 80:
                stu["final_grade"] = "B"
            elif total_score >= 70:
                stu["final_grade"] = "C"
            elif total_score >= 60:
                stu["final_grade"] = "D"
            else:
                stu["final_grade"] = "F"

            stu["remarks"] = "Passed" if stu["final_grade"] != "F" else "Failed"

        # --- DB insert ---
        for stu in students:
            if not any(isinstance(v, (str, int, float)) and str(v).strip() for v in stu.values()):
                continue
            try:
                # NEW: Insert school + batch into immersion_batches if not exists
                school = get_student_value(stu, "school")
                batch = get_student_value(stu, "batch")

                if school and batch:
                    row = config.fetch_one(
                        "SELECT id FROM immersion_batches WHERE school=%s AND batch=%s",
                        (school, batch)
                    )
                    if not row:
                        config.execute_query(
                            "INSERT INTO immersion_batches (school, batch) VALUES (%s, %s)",
                            (school, batch)
                        )

                # Insert student record
                config.execute_query("""
                    INSERT INTO immersion_records (
                        last_name, first_name, middle_name, strand, department,
                        WI, CO, 5S, BO, CBO, SDG,
                        OHSA, WE, UJC, ISO, PO, HR,
                        PERDEV, SUPP, DS,
                        total_score, written_rating, performance_rating, final_grade, remarks
                    )
                    VALUES (
                        %s, %s, %s, %s, %s,
                        %s, %s, %s, %s, %s, %s,
                        %s, %s, %s, %s, %s, %s,
                        %s, %s, %s,
                        %s, %s, %s, %s, %s
                    )
                """, (
                    get_student_value(stu, "last_name"),
                    get_student_value(stu, "first_name"),
                    get_student_value(stu, "middle_name"),
                    get_student_value(stu, "strand"),
                    get_student_value(stu, "department"),

                    int(stu["wi"]), int(stu["co"]), int(stu["5s"]), int(stu["bo"]), int(stu["cbo"]), int(stu["sdg"]),
                    int(stu["ohsa"]), int(stu["we"]), int(stu["ujc"]), int(stu["iso"]), int(stu["po"]), int(stu["hr"]),
                    int(stu["perdev"]), int(stu["supp"]), int(stu["ds"]),

                    float(stu["total_score"]), float(stu["written_rating"]), float(stu["performance_rating"]),
                    stu["final_grade"], stu["remarks"]
                ))
            except Exception as e:
                print(f"❌ DB insert failed for student {stu}: {e}")

        # Return file
        force_full_calc_on_load(wb)
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        filename = f"IMMERSION-GENERATED-{datetime.now().strftime('%Y%m%d-%H%M%S')}.xlsx"
        return send_file(
            output,
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


# ---------- Entry ----------

if __name__ == "__main__":
    app.run(debug=True, host='0.0.0.0', port=5000)
